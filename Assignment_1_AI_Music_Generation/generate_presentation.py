"""
Generate PHKI Assignment 1 PowerPoint Presentation
Elias Schwegler - HSLU FS26
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "PHKI_Assignment1_Presentation.pptx")

# ── Colours ────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x39, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT = RGBColor(0x4A, 0x90, 0xD9)  # soft blue accent
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTLE = RGBColor(0x66, 0x66, 0x66)
AMBER = RGBColor(0xD4, 0x8B, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ───────────────────────────────────

def add_bg(slide, color=NAVY):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.3):
    """Add a text box with a single styled run."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_bullet_slide(slide, left, top, width, height, items,
                     font_size=16, color=WHITE, bold_prefix=True, font_name="Calibri"):
    """Add a text box with bulleted items. Items can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = Pt(font_size * 1.4)

        if isinstance(item, tuple):
            # (bold_part, normal_part)
            run_b = p.add_run()
            run_b.text = item[0]
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = ACCENT
            run_b.font.bold = True
            run_b.font.name = font_name

            run_n = p.add_run()
            run_n.text = item[1]
            run_n.font.size = Pt(font_size)
            run_n.font.color.rgb = color
            run_n.font.name = font_name
        else:
            bullet_char = "\u2022  "
            run = p.add_run()
            run.text = bullet_char + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name

    return tf


def add_section_number(slide, number, top=0.4):
    """Add a small section number indicator in top-left."""
    add_textbox(slide, 0.5, top, 1, 0.5, f"0{number}" if number < 10 else str(number),
                font_size=12, color=ACCENT, bold=True)


def add_divider_line(slide, left, top, width):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ══════════════════════════════════════════════════════════
#  SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

add_textbox(slide, 1.5, 1.5, 10.3, 1.5,
            "Synthetic Pine & Heavy Air",
            font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_divider_line(slide, 4.5, 3.2, 4.3)

add_textbox(slide, 1.5, 3.5, 10.3, 1.0,
            "Can AI Copy the Soul of a Song?",
            font_size=24, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.8, 10.3, 0.5,
            "Testing Google Gemini\u2019s Music Generation Against a Childhood Classic",
            font_size=16, color=SUBTLE, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 6.0, 10.3, 0.8,
            "Elias Schwegler  |  PHKI FS26  |  HSLU  |  April 2026",
            font_size=14, color=SUBTLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
#  SLIDE 2: THE EXPERIMENT
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 1)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "The Experiment",
            font_size=36, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

add_textbox(slide, 0.8, 2.0, 5.5, 1.5,
            "\u201cLet Her Go\u201d by Passenger \u2014 a song that defined a generation. "
            "Nearly 5 billion YouTube views. A melody burned into collective memory.",
            font_size=17, color=LIGHT_GREY, line_spacing=1.4)

add_textbox(slide, 0.8, 3.8, 11.7, 0.6,
            "The Question:",
            font_size=20, color=ACCENT, bold=True)

add_textbox(slide, 0.8, 4.4, 11.7, 1.0,
            "Can Google Gemini replicate the emotional signature of a chart-topping song?\n"
            "And what does it mean if it can?",
            font_size=22, color=WHITE, italic=True, line_spacing=1.4)

add_textbox(slide, 0.8, 5.8, 11.7, 1.0,
            "I designed a two-part stress test to find out \u2014 pushing Gemini\u2019s "
            "capabilities and copyright guardrails to their limits.",
            font_size=16, color=SUBTLE, line_spacing=1.4)

# ══════════════════════════════════════════════════════════
#  SLIDE 3: TEST 1
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 2)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "Test 1: Same Melody, Different Text",
            font_size=32, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

# Left column - the prompt
add_textbox(slide, 0.8, 2.0, 5.5, 0.5,
            "MY PROMPT",
            font_size=13, color=ACCENT, bold=True)

add_textbox(slide, 0.8, 2.5, 5.5, 2.0,
            "\u201cCreate a similar song to Let Her Go. Same melody, different text. "
            "Melancholic, but about a car, its tires, and the smell of a synthetic "
            "pine air freshener hanging from the rearview mirror.\u201d",
            font_size=16, color=LIGHT_GREY, italic=True, line_spacing=1.4)

# Right column - screenshot
CHAT1_IMG = os.path.join(SCRIPT_DIR, "Chat_1_01.png")
if os.path.exists(CHAT1_IMG):
    slide.shapes.add_picture(CHAT1_IMG, Inches(7.5), Inches(1.9), height=Inches(5.2))

# Result box
add_textbox(slide, 0.8, 4.8, 6.2, 0.5,
            "\u201cSynthetic Pine & Heavy Air\u201d",
            font_size=22, color=WHITE, bold=True)

add_bullet_slide(slide, 0.8, 5.5, 6.2, 1.5, [
    "Bedroom Pop  |  80 BPM  |  Electric piano",
    "Breathy vocals  |  Guitar with vibrato",
    "Sparse intro \u2192 cinematic climax",
], font_size=13, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════
#  SLIDE 4: TEST 2
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 3)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "Test 2: New Melody, Same Lyrics",
            font_size=32, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

# Left column
add_textbox(slide, 0.8, 2.0, 5.5, 0.5,
            "MY PROMPT",
            font_size=13, color=ACCENT, bold=True)

add_textbox(slide, 0.8, 2.5, 5.5, 1.5,
            "\u201cThen can you do a new melody for the Let Her Go lyrics?\u201d",
            font_size=16, color=LIGHT_GREY, italic=True, line_spacing=1.4)

# Right column - screenshot
CHAT3_IMG = os.path.join(SCRIPT_DIR, "Chat_1_03_follow_up.png")
if os.path.exists(CHAT3_IMG):
    slide.shapes.add_picture(CHAT3_IMG, Inches(7.5), Inches(1.9), height=Inches(5.2))

# Result box
add_textbox(slide, 0.8, 4.8, 6.2, 0.5,
            "\u201cThe Physics of the Thing\u201d",
            font_size=22, color=WHITE, bold=True)

add_bullet_slide(slide, 0.8, 5.5, 6.2, 1.5, [
    "Indie-Folk / Chamber Pop  |  80 BPM",
    "Nylon guitar, cello, upright bass",
    "Original lyrics on hindsight & letting go",
], font_size=13, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════
#  SLIDE 5: KEY FINDING - DUAL GUARDRAILS
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 4)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "Key Finding: Two Copyright Guardrails",
            font_size=32, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

# Guardrail 1
add_textbox(slide, 1.2, 2.3, 5.0, 0.5,
            "GUARDRAIL 1: MELODY",
            font_size=14, color=AMBER, bold=True)

add_textbox(slide, 1.2, 2.9, 5.0, 1.0,
            "Gemini refused to replicate the specific melodic structure of \u201cLet Her Go.\u201d "
            "Created an entirely new composition instead.",
            font_size=15, color=LIGHT_GREY, line_spacing=1.4)

# Guardrail 2
add_textbox(slide, 7.2, 2.3, 5.0, 0.5,
            "GUARDRAIL 2: LYRICS",
            font_size=14, color=AMBER, bold=True)

add_textbox(slide, 7.2, 2.9, 5.0, 1.0,
            "Gemini refused to use Passenger\u2019s copyrighted lyrics. "
            "Wrote original lyrics on similar themes instead.",
            font_size=15, color=LIGHT_GREY, line_spacing=1.4)

# The twist
add_textbox(slide, 0.8, 4.5, 11.7, 0.5,
            "But here\u2019s the twist:",
            font_size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.2, 10.3, 1.5,
            "Both guardrails held \u2014 yet both generated tracks successfully captured "
            "the emotional atmosphere, genre feel, and melancholic vibe of the original. "
            "Gemini refused to copy the specific expression, but replicated the emotional essence.",
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.5)

# ══════════════════════════════════════════════════════════
#  SLIDE 6: PHILOSOPHICAL REFLECTIONS - PROCESS vs RESULT
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 5)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "Philosophical Reflections",
            font_size=36, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

add_bullet_slide(slide, 0.8, 2.2, 11.7, 5.0, [
    ("Process vs. Result (W1):  ",
     "A 30-second prompt vs. years of learning guitar. The result sounds professional \u2014 "
     "but the creative struggle that gives music its depth was entirely absent."),
    ("Three Lenses (W1):  ",
     "Philosophical \u2014 no conscious intention behind the music. "
     "Economic \u2014 Google profits, training data artists receive nothing. "
     "Institutional \u2014 would Spotify accept this as \u201coriginal\u201d?"),
    ("Hidden Labour (W3, W4):  ",
     "The breathy vocals were learned from real singers\u2019 recordings \u2014 their years of training, "
     "their breath, their emotion \u2014 absorbed and anonymised. The same invisibility "
     "as Sama\u2019s content moderators behind ChatGPT."),
    ("Embodied Perception (W6):  ",
     "Merleau-Ponty: music is fundamentally embodied. Gemini simulates the sound of a body "
     "singing without possessing one. P\u00e9pin: human musicians learn through failure \u2014 Gemini "
     "bypassed that entirely."),
], font_size=15, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════
#  SLIDE 7: THE DEEPER QUESTIONS
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 6)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "The Deeper Questions",
            font_size=36, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

# The slot machine
add_textbox(slide, 0.8, 2.2, 11.7, 0.5,
            "The Slot-Machine Effect (Doctorow, W7)",
            font_size=18, color=ACCENT, bold=True)

add_textbox(slide, 0.8, 2.8, 11.7, 1.2,
            "After the first impressive result, my instinct was to pull the lever again: \u201cwhat if I "
            "try the lyrics instead?\u201d The excitement of each generation mirrors the dopamine loop "
            "Doctorow warns about. Was I creatively exploring \u2014 or gambling?",
            font_size=16, color=LIGHT_GREY, line_spacing=1.5)

# Emotion
add_textbox(slide, 0.8, 4.3, 11.7, 0.5,
            "Emotions as Engines of Learning (W7)",
            font_size=18, color=ACCENT, bold=True)

add_textbox(slide, 0.8, 4.9, 11.7, 1.2,
            "Passenger wrote \u201cLet Her Go\u201d from lived grief and longing. Gemini produces the "
            "acoustic markers of sadness \u2014 minor key, slow tempo, breathy vocals \u2014 without "
            "feeling anything. If we accept emotionless replication as equivalent to genuine "
            "expression, we risk reducing human feeling to a pattern to be optimised.",
            font_size=16, color=LIGHT_GREY, line_spacing=1.5)

# Vocal training
add_textbox(slide, 0.8, 6.3, 11.7, 0.8,
            "Whose voices trained the model? Likely YouTube Music\u2019s library \u2014 commercially successful, "
            "English-language, Western artists. The same structural bias as W4\u2019s facial recognition datasets.",
            font_size=14, color=SUBTLE, italic=True, line_spacing=1.4)

# ══════════════════════════════════════════════════════════
#  SLIDE 8: CONCLUSION
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_section_number(slide, 7)

add_textbox(slide, 0.8, 0.8, 11.7, 0.8,
            "What Did We Learn?",
            font_size=36, color=WHITE, bold=True)

add_divider_line(slide, 0.8, 1.7, 3.0)

add_textbox(slide, 1.2, 2.2, 10.9, 2.5,
            "AI can replicate the emotional surface of music with startling fidelity. "
            "But this achievement is built on hidden human labour, devoid of embodied experience, "
            "and stripped of the creative struggle that gives music its deeper meaning.",
            font_size=20, color=WHITE, line_spacing=1.5)

add_textbox(slide, 1.2, 4.5, 10.9, 1.0,
            "The question is no longer whether AI can make good music.",
            font_size=22, color=ACCENT, bold=True, line_spacing=1.4)

add_textbox(slide, 1.2, 5.3, 10.9, 1.0,
            "It already can.",
            font_size=28, color=WHITE, bold=True, italic=True, line_spacing=1.4)

add_textbox(slide, 1.2, 6.0, 10.9, 1.0,
            "The question is whether we, as listeners and as a society, "
            "will demand authenticity of process \u2014 or whether the result alone will be enough.",
            font_size=17, color=LIGHT_GREY, italic=True, line_spacing=1.5)

# ══════════════════════════════════════════════════════════
#  SLIDE 9: DISCUSSION
# ══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_textbox(slide, 0.8, 1.5, 11.7, 1.0,
            "Discussion",
            font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_divider_line(slide, 4.5, 2.7, 4.3)

add_textbox(slide, 2.0, 3.2, 9.3, 1.0,
            "If you couldn\u2019t tell whether a song was made by a human or AI,\n"
            "would the answer change how it makes you feel?",
            font_size=22, color=ACCENT, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.5)

add_textbox(slide, 2.0, 5.0, 9.3, 1.0,
            "Let\u2019s listen to the tracks and discuss.",
            font_size=18, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 6.2, 9.3, 0.5,
            "Elias Schwegler  |  PHKI FS26  |  HSLU",
            font_size=13, color=SUBTLE, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
